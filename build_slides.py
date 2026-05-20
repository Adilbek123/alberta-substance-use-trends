"""Build a GoA-style briefing deck with conversational what/why/so-what bullets.

Structure mirrors a Cabinet/MDM deck section labels:
  1. Cover
  2. Issue
  3. Purpose
  4. Background
  5. Research Question
  6. Approach
  7. Findings
  8. What This Means
  9. Limitations and Next Steps

Black and white, no jargon in bodies. Each bullet structured as
"what / why / so what" in full sentences. Substance reviewed by an
econometrist; identification claims framed honestly (interpretive choice
versus statistical claim).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "slides.pptx")
CHART = os.path.join(HERE, "output", "rdit_opioid_deaths.png")

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0x99, 0x99, 0x99)

FONT = "Calibri"


def title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = BLACK


def hline(slide, y=1.05):
    line = slide.shapes.add_connector(1, Inches(0.6), Inches(y), Inches(12.7), Inches(y))
    line.line.color.rgb = BLACK
    line.line.width = Pt(0.75)


def bullets(slide, left, top, width, height, items, size=14, sub_size=11):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            main, subs = item
        else:
            main, subs = item, []
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = 0
        r = p.add_run()
        r.text = "•  " + main
        r.font.size = Pt(size)
        r.font.name = FONT
        r.font.color.rgb = DARK
        for sub in subs:
            p = tf.add_paragraph()
            p.level = 1
            r = p.add_run()
            r.text = "–  " + sub
            r.font.size = Pt(sub_size)
            r.font.name = FONT
            r.font.color.rgb = GREY
    return tf


def page_num(slide, n):
    tb = slide.shapes.add_textbox(Inches(12.6), Inches(7.05), Inches(0.6), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(n)
    r.font.size = Pt(10)
    r.font.name = FONT
    r.font.color.rgb = LIGHT


def footer_left(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(11), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9)
    r.font.name = FONT
    r.font.color.rgb = LIGHT


blank = prs.slide_layouts[6]
FOOTER = "Alberta opioid mortality around COVID-19  |  independent project"

# ===========================================================================
# Slide 1 - Cover
# ===========================================================================
s = prs.slides.add_slide(blank)

tb = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.2))
tf = tb.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Alberta Opioid Mortality Around the COVID-19 Onset"
r.font.size = Pt(34)
r.font.bold = True
r.font.name = FONT
r.font.color.rgb = BLACK

tb = s.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Interrupted Time Series Analysis Using Federal Data"
r.font.size = Pt(18)
r.font.name = FONT
r.font.color.rgb = GREY

tb = s.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.6))
tf = tb.text_frame
for line, size, color in [
    ("Adilbek Sultanov", 14, BLACK),
    ("Independent project, May 2026", 12, GREY),
    ("github.com/Adilbek123/alberta-substance-use-trends", 11, GREY),
]:
    if tf.text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    r = p.add_run()
    r.text = line
    r.font.size = Pt(size)
    r.font.name = FONT
    r.font.color.rgb = color

page_num(s, 1)

# ===========================================================================
# Slide 2 - Issue
# ===========================================================================
s = prs.slides.add_slide(blank)
title(s, "Issue")
hline(s)

bullets(s, 0.7, 1.5, 12, 5.0, [
    ("Opioid mortality in Alberta has risen sharply since the start of the COVID-19 pandemic, but the size and timing of the shift have not been quantified for this analysis at quarterly resolution.",
     []),
    ("Without a defensible measurement of the post-COVID baseline, performance under the Alberta Recovery Model cannot be evaluated against a meaningful counterfactual.",
     []),
    ("This analysis quantifies the level shift at the COVID onset using publicly available federal data, and tests whether the shift is statistically meaningful or could be explained by random variation in a noisy series.",
     []),
], size=15, sub_size=12)

footer_left(s, FOOTER)
page_num(s, 2)

# ===========================================================================
# Slide 3 - Purpose
# ===========================================================================
s = prs.slides.add_slide(blank)
title(s, "Purpose")
hline(s)

bullets(s, 0.7, 1.5, 12, 5.0, [
    ("Establish whether Alberta's opioid death rate shifted at the COVID-19 onset",
     ["Identify the size of the shift, if any",
      "Test whether the shift is meaningfully larger than the natural noise in the data"]),
    ("Inform the post-COVID baseline used for Recovery Model performance measurement",
     ["The baseline before the pandemic is materially different from the baseline after",
      "Any performance framing needs to be explicit about which baseline it compares to"]),
    ("Demonstrate an analytical approach the Ministry's Business Intelligence team could apply to other indicators",
     ["The same method applies to hospitalization, emergency department, and EMS data",
      "Same data source, same code structure, different outcome variable"]),
], size=14, sub_size=11)

footer_left(s, FOOTER)
page_num(s, 3)

# ===========================================================================
# Slide 4 - Background
# ===========================================================================
s = prs.slides.add_slide(blank)
title(s, "Background")
hline(s)

bullets(s, 0.7, 1.5, 12, 5.0, [
    ("Alberta has been responding to an opioid crisis since the mid-2010s",
     ["Fentanyl emerged in Alberta's unregulated drug supply through 2014-2016, replacing earlier prescription-opioid harms",
      "In 2016 Alberta recorded 611 apparent opioid-related deaths",
      "In May 2017 the province established a Minister's Opioid Emergency Response Commission under the Public Health Act"]),
    ("By 2016-2019, Alberta's opioid mortality had stabilised at a high level",
     ["Quarterly opioid death rate sat around 4 per 100,000 residents from 2016 through early 2020",
      "Roughly 170-200 deaths per quarter, with no sustained upward or downward trend over the four-year period",
      "This is the pre-COVID baseline used in this analysis"]),
    ("In March 2020 the COVID-19 pandemic intersected with the existing opioid crisis",
     ["WHO characterized COVID-19 as a pandemic on 11 March 2020",
      "Alberta declared a provincial public health emergency on 17 March 2020",
      "Federal data shows opioid mortality rose sharply across Canada from 2020 onward; the size and timing of Alberta's specific shift have not been quantified at this resolution in publicly available analysis"]),
], size=13, sub_size=11)

footer_left(s, FOOTER)
page_num(s, 4)

# ===========================================================================
# Slide 5 - Research Question
# ===========================================================================
s = prs.slides.add_slide(blank)
title(s, "Research Question")
hline(s)

bullets(s, 0.7, 1.5, 12, 5.0, [
    ("Did Alberta's opioid death rate shift at the COVID-19 cutoff (2020 Q2) in a way the pre-COVID period cannot explain?",
     ["This is a 'break' question, not a 'cause' question: we ask whether the post-cutoff data are inconsistent with the pre-period pattern",
      "The pre-COVID period gives us a model of how the rate behaved (level, trend, seasonality); extrapolating that model forward gives a counterfactual for what we would expect post-cutoff if nothing had changed",
      "If observed post-cutoff data depart materially from that extrapolation, the shift is real"]),
    ("If a shift exists, two sub-questions follow",
     ["How large is it, in deaths per 100,000 per quarter, and how large is it relative to pre-period variation?",
      "Could a pattern of this size arise from normal fluctuation in a noisy series? This is a precision question, not yet a policy question"]),
], size=14, sub_size=11)

footer_left(s, FOOTER)
page_num(s, 5)

# ===========================================================================
# Slide 6 - Approach
# ===========================================================================
s = prs.slides.add_slide(blank)
title(s, "Approach")
hline(s)

bullets(s, 0.7, 1.4, 12, 5.2, [
    ("To answer this research question I used interrupted time series (segmented regression)",
     ["The method directly tests for a level shift at a known cutoff, with the pre-period model (intercept, slope, seasonality) extrapolated forward as the counterfactual",
      "It estimates three things separately: the existing trajectory before the cutoff, the level gap at the cutoff, and any change in trajectory after",
      "Standard tool in health policy and public health epidemiology for exactly this kind of pre/post comparison"]),
    ("I used federal public data because the analysis must be independent and reproducible",
     ["Quarterly opioid deaths from the Public Health Agency of Canada, 2016 Q1 to 2025 Q3",
      "Quarterly population from Statistics Canada, used to convert counts to a rate per 100,000 so the result is not driven by Alberta's population growth (about 12 per cent over the period)",
      "No internal Government of Alberta data was used"]),
    ("I added four robustness checks because a single estimate from one specification is not enough to defend a finding",
     ["Placebo cutoffs in the pre-period: re-running the model with fake cutoffs at pre-COVID dates generates a reference distribution for the size of 'jumps' the method finds when nothing has happened; the real cutoff should lie outside that distribution",
      "HAC lag sensitivity: re-estimate the standard errors with different lag lengths so the inference does not depend on a single tuning choice",
      "Donut: drop the partial-exposure quarter (2020 Q1) and refit, to confirm the estimate is not driven by one transitional observation",
      "Negative binomial cross-check with population offset: tests both the linear-additive functional form and the OLS Gaussian distributional assumption against a count model better matched to mortality data"]),
], size=13, sub_size=10)

footer_left(s, FOOTER)
page_num(s, 6)

# ===========================================================================
# Slide 7 - Findings
# ===========================================================================
s = prs.slides.add_slide(blank)
title(s, "Findings")
hline(s)

if os.path.exists(CHART):
    s.shapes.add_picture(CHART, Inches(0.5), Inches(1.25), width=Inches(8.4))

bullets(s, 9.2, 1.4, 4.0, 5.5, [
    ("Pre-COVID baseline was flat",
     ["Quarterly rate sat around 4 per 100,000 from 2016 through early 2020",
      "Pre-COVID slope essentially zero (p = 0.94)"]),
    ("Level shift at 2020 Q2",
     ["About +5.5 per 100,000 per quarter (95% CI +3.6 to +7.4)",
      "Equivalent to a rate ratio of 2.4 (count-model cross-check)",
      "The rate roughly doubled (~4 per 100k to ~9.5 immediately post-cutoff)",
      "At Alberta's mean population over the period, this is ~246 additional deaths per quarter, evaluated at the cutoff"]),
    ("Inference is strong",
     ["p < 0.001 in the main spec",
      "Placebo cutoffs in the pre-period maxed out at about one third of the real estimate, in the opposite direction; the real shift is far outside the placebo distribution"]),
    ("Robustness checks all hold",
     ["HAC lag sensitivity: estimate stable across lag choices",
      "Donut: drops the transitional quarter, estimate moves <3%",
      "NB cross-check confirms the rate-ratio interpretation"]),
], size=11, sub_size=9)

footer_left(s, FOOTER)
page_num(s, 7)

# ===========================================================================
# Slide 8 - What This Means
# ===========================================================================
s = prs.slides.add_slide(blank)
title(s, "What This Means")
hline(s)

bullets(s, 0.7, 1.4, 12, 5.2, [
    ("The analysis identifies the level shift at 2020 Q2 that the pre-period trend and seasonality cannot explain",
     ["This is what the statistics support directly; everything below is interpretation"]),
    ("I interpret that shift as the total joint effect of the COVID-19 pandemic bundle",
     ["'Pandemic' here means everything that came together at the cutoff: virus circulation, the public health emergency response, drug supply chain disruption, harm reduction service capacity, mental health and addiction service capacity, isolation, and the economic shock",
      "These are downstream consequences of the pandemic; they did not align by coincidence",
      "This is an interpretive choice, not a statistical claim. The design cannot separate the channels, and the channels cannot meaningfully be removed from 'COVID' without losing the concept itself"]),
    ("The analysis does NOT identify which channel caused most of the damage",
     ["That is a different question (mechanism, not break) and needs different data: drug-checking composition, service utilization records, fentanyl share of toxicology, and so on",
      "Comparing Alberta to other provinces would not solve this: every Canadian province faced COVID and its policy response at the same time, so peer comparison identifies Alberta-specific deviation, not a virus-only channel"]),
    ("The post-cutoff series has been declining since 2024",
     ["This analysis does not test that decline formally; doing so is a natural next step",
      "The Recovery Model's performance framing should be explicit about which baseline it compares to: pre-COVID, post-cutoff peak, or current"]),
], size=12, sub_size=10)

footer_left(s, FOOTER)
page_num(s, 8)

# ===========================================================================
# Slide 9 - Limitations and Next Steps
# ===========================================================================
s = prs.slides.add_slide(blank)
title(s, "Limitations and Next Steps")
hline(s)

bullets(s, 0.7, 1.4, 12, 5.2, [
    ("Limitations to keep in mind",
     ["Single province, single outcome (opioid deaths) - no decomposition by age, sex, zone, or substance type",
      "Small sample for asymptotic inference (n = 39 quarters); n is at the lower end of where Newey-West asymptotics are reasonable, which is why placebos serve as a partial substitute",
      "The cutoff (2020 Q2) was chosen because the emergency was declared 17 March 2020; 2020 Q1 is partial-exposure, the donut robustness confirms the choice does not drive the result",
      "Beta-2 is the level shift conditional on a linear post-trend; the post-period is non-monotonic (peak in 2021 Q4, decline since 2024), so a single level shift does not capture the full dynamic",
      "Seasonality is modelled as additive quarter fixed effects, assumed constant across years (likely fine but not formally tested)",
      "Reporting lag and revisions: the most recent quarters may move when PHAC updates the data"]),
    ("What I would do next",
     ["Mechanism decomposition: bring in supply-side, service-capacity, and outcome data to attribute share of the shift to each channel",
      "Heterogeneity: by age, sex, zone (Edmonton, Calgary, rural), and substance type (fentanyl, carfentanil, other)",
      "Post-period dynamics: model the rise-to-peak-and-decline shape directly, rather than collapsing it into a single level shift",
      "Cumulative excess deaths through end-of-sample as a policy-relevant total"]),
    ("Reproducibility",
     ["Code, data, and full documentation: github.com/Adilbek123/alberta-substance-use-trends"]),
], size=12, sub_size=10)

footer_left(s, FOOTER)
page_num(s, 9)

prs.save(OUT)
print(f"Saved: {OUT}")
